# streamlit_app.py
"""
Storybook / Magazine-style Weekly Report Builder for JUNA PV
- Manual inputs + uploads (CSV, images)
- Magazine-style PDF export (multi-page) using matplotlib PdfPages
- PPTX export using python-pptx
"""

import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from io import BytesIO
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import textwrap

st.set_page_config(page_title="Storybook — Weekly Report (Magazine)", layout="wide")

# ---------- Helper functions ----------
def to_float_safe(x):
    try:
        return float(x)
    except:
        return None

def load_csv_uploader(uploader):
    if uploader is None:
        return None
    try:
        df = pd.read_csv(uploader)
        return df
    except Exception as e:
        try:
            df = pd.read_excel(uploader)
            return df
        except Exception:
            st.error(f"Failed to read uploaded file: {e}")
            return None

def small_text_box(text, width=70):
    return "\n".join(textwrap.wrap(str(text), width=width))

def draw_kpi_grid(ax, kpis):
    ax.axis('off')
    # Display KPIs as big text blocks
    rows = 2
    cols = 3
    keys = list(kpis.keys())
    for i in range(rows):
        for j in range(cols):
            idx = i*cols + j
            if idx >= len(keys):
                continue
            key = keys[idx]
            val = kpis[key]
            x = j / cols
            y = 1 - (i+1)/rows + 0.02
            ax.text(x + 0.02, y, f"{key}", fontsize=12, weight='bold', transform=ax.transAxes)
            ax.text(x + 0.02, y - 0.07, f"{val}", fontsize=20, transform=ax.transAxes)

def save_imgfile_to_bytes(uploaded_file):
    # return PIL Image and bytes
    try:
        img = Image.open(uploaded_file).convert("RGB")
        bio = BytesIO()
        img.save(bio, format="PNG")
        bio.seek(0)
        return img, bio
    except Exception as e:
        st.error(f"Failed to load image: {e}")
        return None, None

def generate_magazine_pdf_buf(cover_image, logo_image, kpis, highlights, issues, plan, df_ts, df_env, df_break, df_pp, df_robot, photo_list, timeline_events):
    """
    Compose multi-page PDF using matplotlib PdfPages and return BytesIO buffer.
    """
    buf = BytesIO()
    with PdfPages(buf) as pdf:
        # --- Cover page ---
        fig = plt.figure(figsize=(8.27, 11.69))  # A4
        ax = fig.add_axes([0,0,1,1])
        ax.axis('off')
        if cover_image is not None:
            # show cover image as background (fit)
            ax_im = fig.add_axes([0.05,0.3,0.9,0.6])
            ax_im.imshow(cover_image)
            ax_im.axis('off')
        ax.text(0.06, 0.22, "JUNA PV", fontsize=28, weight='bold')
        ax.text(0.06, 0.18, f"Weekly Report — {kpis.get('Week','')}", fontsize=12)
        ax.text(0.06, 0.14, f"Date Range: {kpis.get('Date Range','')}", fontsize=10)
        if logo_image is not None:
            # logo on top-right
            ax_logo = fig.add_axes([0.7,0.75,0.25,0.2])
            ax_logo.imshow(logo_image)
            ax_logo.axis('off')
        pdf.savefig(fig, bbox_inches='tight')
        plt.close(fig)

        # --- Executive summary page ---
        fig = plt.figure(figsize=(8.27, 11.69))
        ax = fig.add_axes([0.05,0.05,0.9,0.9])
        ax.axis('off')
        ax.text(0.02, 0.92, "Executive Summary", fontsize=20, weight='bold')
        draw_kpi_grid(ax, kpis)
        ax.text(0.02, 0.38, "Highlights:", fontsize=12, weight='bold')
        ax.text(0.02, 0.35, small_text_box(highlights or "No highlights provided"), fontsize=10)
        ax.text(0.02, 0.25, "Major Issues:", fontsize=12, weight='bold')
        ax.text(0.02, 0.22, small_text_box(issues or "No issues provided"), fontsize=10)
        ax.text(0.02, 0.12, "Plan for Next Week:", fontsize=12, weight='bold')
        ax.text(0.02, 0.09, small_text_box(plan or "No plan provided"), fontsize=10)
        pdf.savefig(fig, bbox_inches='tight')
        plt.close(fig)

        # --- Performance / Charts page ---
        fig = plt.figure(figsize=(8.27, 11.69))
        gs = fig.add_gridspec(3,1)
        ax1 = fig.add_subplot(gs[0,:])
        ax2 = fig.add_subplot(gs[1,:])
        ax3 = fig.add_subplot(gs[2,:])
        if df_ts is not None and not df_ts.empty:
            try:
                dfp = df_ts.copy()
                if 'date' in dfp.columns:
                    dfp['date'] = pd.to_datetime(dfp['date'])
                    dfp = dfp.sort_values('date')
                    ax1.plot(dfp['date'], dfp.iloc[:,1], marker='o')
                else:
                    ax1.plot(dfp.iloc[:,0], dfp.iloc[:,1], marker='o')
                ax1.set_title("Weekly Active Power")
                ax1.set_ylabel("MW")
            except Exception as e:
                ax1.text(0.5,0.5,f"Failed to plot timeseries: {e}", ha='center')
        else:
            ax1.text(0.5,0.5,"No timeseries provided", ha='center')

        # Environmental if available
        if df_env is not None and not df_env.empty:
            try:
                env = df_env.copy()
                env['date'] = pd.to_datetime(env['date'])
                env = env.sort_values('date')
                if 'irradiance' in [c.lower() for c in env.columns]:
                    col = [c for c in env.columns if 'irradiance' in c.lower()][0]
                    ax2.plot(env['date'], env[col], marker='.', label='Irradiance')
                if 'temperature' in [c.lower() for c in env.columns]:
                    col2 = [c for c in env.columns if 'temperature' in c.lower()][0]
                    ax2.plot(env['date'], env[col2], marker='.', label='Temperature')
                ax2.legend()
                ax2.set_title("Environmental data")
            except Exception as e:
                ax2.text(0.5,0.5,f"Failed to plot env: {e}", ha='center')
        else:
            ax2.text(0.5,0.5,"No environmental data provided", ha='center')

        # small summary table area
        ax3.axis('off')
        text = ""
        if df_break is not None and not df_break.empty:
            try:
                total_downtime = pd.to_numeric(df_break['downtime_minutes'], errors='coerce').sum()
                text += f"Total Downtime (mins): {int(total_downtime)}\n"
            except:
                pass
        if df_robot is not None and not df_robot.empty:
            text += f"Robot logs: {len(df_robot)} entries\n"
        ax3.text(0.02, 0.98, text, verticalalignment='top', fontsize=10)
        pdf.savefig(fig, bbox_inches='tight')
        plt.close(fig)

        # --- Breakdowns page ---
        if df_break is not None and not df_break.empty:
            fig, ax = plt.subplots(figsize=(8.27,11.69))
            ax.axis('off')
            ax.text(0.02, 0.95, "Breakdown Log (sample rows)", fontsize=16, weight='bold')
            # show first 12 rows as text
            rows = df_break.head(12).to_dict(orient='records')
            y = 0.9
            for r in rows:
                line = " | ".join([f"{k}:{v}" for k,v in r.items()])
                ax.text(0.02, y, small_text_box(line, width=120), fontsize=9)
                y -= 0.06
                if y < 0.06:
                    pdf.savefig(fig, bbox_inches='tight')
                    plt.close(fig)
                    fig, ax = plt.subplots(figsize=(8.27,11.69))
                    ax.axis('off')
                    y = 0.95
            pdf.savefig(fig, bbox_inches='tight')
            plt.close(fig)

        # --- Punchpoints page ---
        if df_pp is not None and not df_pp.empty:
            fig, ax = plt.subplots(figsize=(8.27,11.69))
            ax.axis('off')
            ax.text(0.02, 0.95, "Punch Points Summary", fontsize=16, weight='bold')
            rows = df_pp.head(20).to_dict(orient='records')
            y = 0.9
            for r in rows:
                line = " | ".join([f"{k}:{v}" for k,v in r.items()])
                ax.text(0.02, y, small_text_box(line, width=120), fontsize=9)
                y -= 0.045
                if y < 0.06:
                    pdf.savefig(fig, bbox_inches='tight')
                    plt.close(fig)
                    fig, ax = plt.subplots(figsize=(8.27,11.69))
                    ax.axis('off')
                    y = 0.95
            pdf.savefig(fig, bbox_inches='tight')
            plt.close(fig)

        # --- Photo gallery pages (2 images per page) ---
        if photo_list:
            chunks = [photo_list[i:i+2] for i in range(0, len(photo_list), 2)]
            for chunk in chunks:
                fig = plt.figure(figsize=(8.27,11.69))
                for i, img in enumerate(chunk):
                    ax = fig.add_axes([0.05, 0.55 - i*0.5, 0.9, 0.45])
                    ax.imshow(img)
                    ax.axis('off')
                pdf.savefig(fig, bbox_inches='tight')
                plt.close(fig)

        # --- Timeline pages ---
        if timeline_events:
            fig = plt.figure(figsize=(8.27,11.69))
            ax = fig.add_axes([0.05,0.05,0.9,0.9])
            ax.axis('off')
            ax.text(0.02, 0.95, "Timeline of Events", fontsize=18, weight='bold')
            y = 0.9
            for ev in timeline_events:
                date = ev.get('date','')
                title = ev.get('title','')
                desc = ev.get('desc','')
                ax.text(0.02, y, f"{date} — {title}", fontsize=12, weight='bold')
                y -= 0.03
                ax.text(0.04, y, small_text_box(desc, width=120), fontsize=10)
                y -= 0.06
                if y < 0.06:
                    pdf.savefig(fig, bbox_inches='tight')
                    plt.close(fig)
                    fig = plt.figure(figsize=(8.27,11.69))
                    ax = fig.add_axes([0.05,0.05,0.9,0.9])
                    ax.axis('off')
                    y = 0.9
            pdf.savefig(fig, bbox_inches='tight')
            plt.close(fig)

    buf.seek(0)
    return buf

def generate_pptx_buf(cover_image_bytes, logo_bytes, kpis, highlights, issues, plan, df_ts, df_env, df_break, df_pp, df_robot, photo_bytes_list, timeline_events):
    prs = Presentation()
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "JUNA PV — Weekly Report"
    try:
        slide.placeholders[1].text = f"Week: {kpis.get('Week','')}\nDate Range: {kpis.get('Date Range','')}"
    except:
        pass
    # KPIs slide
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Executive Summary"
    tf = s.shapes.add_textbox(Inches(0.4), Inches(1.2), Inches(9), Inches(3)).text_frame
    tf.clear()
    for k,v in kpis.items():
        p = tf.add_paragraph()
        p.text = f"{k}: {v}"
    # Highlights slide
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Highlights"
    tx = s2.shapes.placeholders[1].text_frame
    tx.clear()
    tx.text = highlights or "No highlights provided"
    # Add a sample chart slide (timeseries) as an image
    img_stream = BytesIO()
    if df_ts is not None and not df_ts.empty:
        try:
            dfp = df_ts.copy()
            dfp['date'] = pd.to_datetime(dfp['date'])
            fig, ax = plt.subplots(figsize=(10,3))
            ax.plot(dfp['date'], dfp.iloc[:,1])
            ax.set_title("Weekly Active Power")
            fig.tight_layout()
            fig.savefig(img_stream, format='png')
            plt.close(fig)
            img_stream.seek(0)
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "Performance"
            slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.6), width=Inches(9))
        except Exception:
            pass
    # Photo slides (one per photo)
    for pb in photo_bytes_list:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Photo"
        try:
            slide.shapes.add_picture(pb, Inches(0.5), Inches(1.2), width=Inches(9))
        except Exception:
            pass
    # Timeline slide
    if timeline_events:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Timeline"
        tx = slide.shapes.add_textbox(Inches(0.4), Inches(1.2), Inches(9), Inches(4)).text_frame
        tx.clear()
        for ev in timeline_events:
            p = tx.add_paragraph()
            p.text = f"{ev.get('date','')}: {ev.get('title','')} - {ev.get('desc','')}"
    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out

# ---------- UI ----------

st.title("Storybook / Magazine-style Weekly Report Builder — JUNA PV")
st.markdown("Build a beautiful weekly magazine-style report by uploading data and photos. Exports: Magazine PDF & PPTX.")

# HEADER inputs
with st.expander("Header (cover & basic info)", expanded=True):
    col1, col2 = st.columns([3,1])
    with col1:
        plant_name = st.text_input("Plant name", value="JUNA PV")
        week_label = st.text_input("Week label", value="Week 47 — 20 Nov 2025 to 26 Nov 2025")
        date_range = st.text_input("Date range", value="20 Nov 2025 - 26 Nov 2025")
    with col2:
        logo_file = st.file_uploader("Upload logo (png/jpg)", type=["png","jpg","jpeg"], key="logo")
        cover_file = st.file_uploader("Upload cover image (optional)", type=["png","jpg","jpeg"], key="cover")

# KPI inputs
with st.expander("KPIs (enter manually)", expanded=True):
    col1, col2, col3 = st.columns(3)
    with col1:
        weekly_energy = st.text_input("Weekly Energy (GWh)", value="9.93")
        mtd_energy = st.text_input("MTD Energy (GWh)", value="25.19")
    with col2:
        ytd_energy = st.text_input("YTD Energy (GWh)", value="206.41")
        plant_avail = st.text_input("Plant Availability (%)", value="100")
    with col3:
        curtailment = st.text_input("Curtailment (%)", value="0")
        pr = st.text_input("PR (%)", value="0")

kpis = {
    "Plant": plant_name,
    "Week": week_label,
    "Date Range": date_range,
    "Weekly Energy (GWh)": weekly_energy,
    "MTD Energy (GWh)": mtd_energy,
    "YTD Energy (GWh)": ytd_energy,
    "Plant Availability (%)": plant_avail,
    "Curtailment (%)": curtailment,
    "PR (%)": pr
}

# Timeseries uploads
with st.expander("Timeseries / Environmental data", expanded=False):
    ts_file = st.file_uploader("Upload timeseries CSV (date, value) for active power", type=["csv","xlsx"], key="ts")
    env_file = st.file_uploader("Upload environmental CSV (date, irradiance, temperature optional)", type=["csv","xlsx"], key="env")
    df_ts = load_csv_uploader(ts_file) if ts_file else None
    df_env = load_csv_uploader(env_file) if env_file else None
    if df_ts is not None:
        st.success("Timeseries loaded — first rows:")
        st.dataframe(df_ts.head())

# Highlights / notes
with st.expander("Write Highlights, Issues, Plan", expanded=True):
    highlights = st.text_area("Highlights (bullet points)", value="- Thermography completed\n- Robot trials conducted")
    issues = st.text_area("Major Issues / Incidents", value="- Cable theft in PB-20")
    plan = st.text_area("Plan for next week", value="- Continue robot deployment and SCADA updates")

# Breakdowns upload
with st.expander("Breakdown log (CSV)", expanded=False):
    break_file = st.file_uploader("Upload breakdown CSV (start_date,end_date,block,inverter,unit,fault,rectification,downtime_minutes)", type=["csv","xlsx"], key="break")
    df_break = load_csv_uploader(break_file) if break_file else None
    if df_break is not None:
        st.dataframe(df_break.head())

# Punchpoints upload
with st.expander("Punch Points (CSV)", expanded=False):
    pp_file = st.file_uploader("Upload punchpoints CSV (block,raised,closed,pending)", type=["csv","xlsx"], key="pp")
    df_pp = load_csv_uploader(pp_file) if pp_file else None
    if df_pp is not None:
        st.dataframe(df_pp.head())

# Robot cleaning logs
with st.expander("Robot / Module Cleaning logs (CSV)", expanded=False):
    robot_file = st.file_uploader("Upload robot/cleaning CSV (date,block,status)", type=["csv","xlsx"], key="robot")
    df_robot = load_csv_uploader(robot_file) if robot_file else None
    if df_robot is not None:
        st.dataframe(df_robot.head())

# Photo gallery
with st.expander("Photo gallery (upload multiple images)", expanded=True):
    photos = st.file_uploader("Upload photos (jpg/png) — maintenance, thermography, drone, etc.", type=["png","jpg","jpeg"], accept_multiple_files=True)
    photo_imgs = []
    photo_bytes = []
    if photos:
        cols = st.columns(3)
        for i, p in enumerate(photos):
            img, bio = save_imgfile_to_bytes(p)
            if img is not None:
                photo_imgs.append(img)
                photo_bytes.append(bio.getvalue())
                with cols[i % 3]:
                    st.image(img, caption=p.name, use_column_width=True)

# Timeline builder (manual entries)
with st.expander("Timeline of events (add entries)", expanded=True):
    st.info("Add events to the timeline. You can attach one photo per event (optional).")
    if 'timeline_events' not in st.session_state:
        st.session_state['timeline_events'] = []
    col1, col2, col3 = st.columns([2,4,1])
    with col1:
        te_date = st.date_input("Event date", value=datetime.today())
    with col2:
        te_title = st.text_input("Event title", "")
    with col3:
        te_file = st.file_uploader("Event photo (optional)", type=["png","jpg","jpeg"], key=f"tev{len(st.session_state['timeline_events'])}")
    te_desc = st.text_area("Description", "")
    if st.button("Add timeline event"):
        ev = {"date": te_date.strftime("%Y-%m-%d"), "title": te_title, "desc": te_desc}
        if te_file:
            img, bio = save_imgfile_to_bytes(te_file)
            if img is not None:
                ev['_img'] = img
                ev['_img_bytes'] = bio.getvalue()
        st.session_state['timeline_events'].append(ev)
        st.success("Event added")
    if st.session_state['timeline_events']:
        st.write("Current timeline events:")
        for ev in st.session_state['timeline_events']:
            st.markdown(f"- **{ev['date']}** — {ev['title']} — {ev['desc']}")

timeline_events = st.session_state.get('timeline_events', [])

st.markdown("---")
st.header("Preview & Export")

# Preview panel (simple)
st.subheader("Preview Executive Summary")
left, right = st.columns([1,2])
with left:
    st.markdown(f"**Plant:** {plant_name}")
    st.markdown(f"**Week:** {week_label}")
    st.markdown(f"**Date Range:** {date_range}")
    st.markdown(f"**Weekly Energy (GWh):** {weekly_energy}")
    st.markdown(f"**MTD Energy (GWh):** {mtd_energy}")
    st.markdown(f"**YTD Energy (GWh):** {ytd_energy}")
with right:
    st.markdown("**Highlights**")
    st.markdown(highlights.replace("\n","  \n"))
    st.markdown("**Issues**")
    st.markdown(issues.replace("\n","  \n"))

# Export buttons
st.subheader("Export")
col1, col2 = st.columns(2)
with col1:
    if st.button("Generate magazine-style PDF"):
        # prepare images
        cover_img = None
        logo_img = None
        if cover_file:
            try:
                cover_img = Image.open(cover_file).convert("RGB")
            except:
                cover_img = None
        if logo_file:
            try:
                logo_img = Image.open(logo_file).convert("RGB")
            except:
                logo_img = None
        pdf_buf = generate_magazine_pdf_buf(cover_img, logo_img, kpis, highlights, issues, plan, df_ts, df_env, df_break, df_pp, df_robot, photo_imgs, timeline_events)
        st.success("PDF Ready — download below")
        st.download_button("Download Magazine PDF", data=pdf_buf.getvalue(), file_name="juna_weekly_magazine.pdf", mime="application/pdf")
with col2:
    if st.button("Generate PPTX (magazine slides)"):
        # prepare bytes for images
        cover_bytes = None
        logo_bytes = None
        if cover_file:
            try:
                cover_bytes = Image.open(cover_file).convert("RGB")
                b = BytesIO(); cover_bytes.save(b, format="PNG"); b.seek(0); cover_bytes = b.getvalue()
            except:
                cover_bytes = None
        if logo_file:
            try:
                logo_bytes = Image.open(logo_file).convert("RGB")
                b = BytesIO(); logo_bytes.save(b, format="PNG"); b.seek(0); logo_bytes = b.getvalue()
            except:
                logo_bytes = None
        photo_bytes_list = photo_bytes  # raw bytes
        pptx_buf = generate_pptx_buf(cover_bytes, logo_bytes, kpis, highlights, issues, plan, df_ts, df_env, df_break, df_pp, df_robot, photo_bytes_list, timeline_events)
        st.success("PPTX ready — download below")
        st.download_button("Download PPTX", data=pptx_buf.getvalue(), file_name="juna_weekly_presentation.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

st.markdown("---")
st.caption("Tip: for best magazine layout, upload a high-res cover image (A4 ratio) and 4-8 photos for the gallery.")
